#!/usr/bin/env python3
"""Generate YogaTracker project presentation (.pptx) with images."""

import importlib.util
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PRIMARY = RGBColor(0x2D, 0x6A, 0x4F)
PRIMARY_LIGHT = RGBColor(0x40, 0x9B, 0x6E)
ACCENT = RGBColor(0x95, 0xD5, 0xB2)
DARK = RGBColor(0x1B, 0x43, 0x32)
MUTED = RGBColor(0x5C, 0x6B, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xFA, 0xF9)

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "docs" / "YogaTracker-Project-Presentation.pptx"
ASSETS = ROOT / "docs" / "presentation-assets"


def load_assets() -> dict[str, Path]:
    spec = importlib.util.spec_from_file_location(
        "gen_assets",
        Path(__file__).parent / "generate-presentation-assets.py",
    )
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    mod.build_all()
    return {
        "logo": ASSETS / "logo.png",
        "session": ASSETS / "session-screen.png",
        "dashboard": ASSETS / "dashboard-screen.png",
        "landing": ASSETS / "landing-screen.png",
        "architecture": ASSETS / "architecture.png",
        "poses": ASSETS / "pose-grid.png",
        "flow": ASSETS / "flow-diagram.png",
    }


def add_image(slide, path: Path, left, top, width, height=None):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height) if height else None)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(9), Inches(0.7))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(9), Inches(0.4))
        stf = sub.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(14)
        sp.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], left=0.7, top=1.85, width=8.8, height=5.0, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_two_column(slide, left_title, left_items, right_title, right_items):
    lt = slide.shapes.add_textbox(Inches(0.55), Inches(1.75), Inches(4.3), Inches(0.4))
    ltp = lt.text_frame.paragraphs[0]
    ltp.text = left_title
    ltp.font.bold = True
    ltp.font.size = Pt(16)
    ltp.font.color.rgb = PRIMARY

    lb = slide.shapes.add_textbox(Inches(0.55), Inches(2.15), Inches(4.3), Inches(4.5))
    ltf = lb.text_frame
    ltf.word_wrap = True
    for i, item in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)

    rt = slide.shapes.add_textbox(Inches(5.15), Inches(1.75), Inches(4.3), Inches(0.4))
    rtp = rt.text_frame.paragraphs[0]
    rtp.text = right_title
    rtp.font.bold = True
    rtp.font.size = Pt(16)
    rtp.font.color.rgb = PRIMARY

    rb = slide.shapes.add_textbox(Inches(5.15), Inches(2.15), Inches(4.3), Inches(4.5))
    rtf = rb.text_frame
    rtf.word_wrap = True
    for i, item in enumerate(right_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def build() -> Path:
    img = load_assets()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Title + session preview
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PRIMARY)
    add_image(slide, img["logo"], 0.9, 0.55, 1.1)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(4.5), Inches(1.0))
    tp = title.text_frame.paragraphs[0]
    tp.text = "YogaTracker"
    tp.font.size = Pt(48)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(4.5), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Fitness & Yoga Form Tracking with Webcam AI"
    sp.font.size = Pt(18)
    sp.font.color.rgb = ACCENT

    tag = slide.shapes.add_textbox(Inches(0.8), Inches(3.55), Inches(4.5), Inches(0.8))
    tgp = tag.text_frame.paragraphs[0]
    tgp.text = "Live pose detection • Rep counting • Progress dashboard"
    tgp.font.size = Pt(12)
    tgp.font.color.rgb = WHITE

    live = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(4.5), Inches(0.4))
    lp = live.text_frame.paragraphs[0]
    lp.text = "yoga-tracker-kappa.vercel.app"
    lp.font.size = Pt(12)
    lp.font.color.rgb = ACCENT

    add_image(slide, img["session"], 5.2, 1.0, 4.5, 5.8)

    # Slide 2 — Overview + landing screenshot
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Project Overview")
    add_bullets(
        slide,
        [
            "Full-stack web app for home workouts with real-time form feedback.",
            "Uses webcam + MediaPipe to track body landmarks in the browser.",
            "Supports fitness reps (squat, push-up, lunge) and yoga holds.",
            "Scores joint angles, counts reps, and gives voice/visual coaching.",
            "Saves sessions, streaks, and charts to a Supabase database.",
            "Deployed on Vercel with GitHub source control.",
        ],
        width=4.6,
        height=4.8,
    )
    add_image(slide, img["landing"], 5.35, 1.75, 4.2, 2.35)
    cap = slide.shapes.add_textbox(Inches(5.35), Inches(4.2), Inches(4.2), Inches(0.35))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Landing page"
    cp.font.size = Pt(11)
    cp.font.color.rgb = MUTED
    cp.alignment = PP_ALIGN.CENTER

    # Slide 3 — Problem & Solution
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Problem & Solution")
    add_two_column(
        slide,
        "Problem",
        [
            "• Hard to check form without a trainer",
            "• No record of reps, time, or accuracy at home",
            "• Generic fitness apps don't score posture live",
            "• Yoga beginners need alignment cues in real time",
        ],
        "Solution",
        [
            "• Browser app — no install, works on laptop/phone (HTTPS)",
            "• Skeleton overlay + form score on the video feed",
            "• Exercise demo panel and coaching tips while you move",
            "• Dashboard with weekly reports and session history",
        ],
    )

    # Slide 4 — Features + app screenshots
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Key Features")
    add_two_column(
        slide,
        "User app",
        [
            "• Sign up / login (Supabase Auth)",
            "• 3-step onboarding (profile, goal, equipment)",
            "• Personalized workout plan",
            "• Live session: camera, exercise picker, demo overlay",
            "• Rep counter for strength moves",
            "• Voice + on-screen feedback",
        ],
        "Backend & security",
        [
            "• PostgreSQL tables + Row Level Security",
            "• REST APIs (sessions, exercises, analytics)",
            "• JWT-protected routes via middleware",
            "• Zod validation on forms and API bodies",
            "• Exercise logs and daily progress logs",
            "• Admin panel for users and poses",
        ],
    )
    add_image(slide, img["session"], 0.55, 5.55, 4.3, 1.55)
    add_image(slide, img["dashboard"], 5.15, 5.55, 4.3, 1.55)

    # Slide 5 — Tech stack
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Technology Stack")
    rows = [
        ("Frontend", "Next.js 16, React 19, TypeScript, Tailwind, shadcn/ui"),
        ("Pose AI", "MediaPipe Tasks Vision (PoseLandmarker)"),
        ("Auth & DB", "Supabase Auth, PostgreSQL, Storage"),
        ("Charts", "Recharts"),
        ("Validation", "Zod + React Hook Form"),
        ("Hosting", "Vercel (production)"),
        ("Repo", "GitHub — yoga-pose-Tracker"),
    ]
    y = 1.9
    for layer, tech in rows:
        layer_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(2.0), Inches(0.35))
        lp = layer_box.text_frame.paragraphs[0]
        lp.text = layer
        lp.font.bold = True
        lp.font.size = Pt(15)
        lp.font.color.rgb = PRIMARY

        tech_box = slide.shapes.add_textbox(Inches(2.7), Inches(y), Inches(6.7), Inches(0.35))
        tp = tech_box.text_frame.paragraphs[0]
        tp.text = tech
        tp.font.size = Pt(15)
        tp.font.color.rgb = DARK
        y += 0.55

    # Slide 6 — Architecture diagram
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "System Architecture")
    add_image(slide, img["architecture"], 0.5, 1.55, 9.0, 5.5)

    # Slide 7 — Session flow + screenshot
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Session Flow")
    add_image(slide, img["flow"], 0.55, 1.65, 8.9, 1.45)
    add_image(slide, img["session"], 0.55, 3.25, 8.9, 3.85)
    cap = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(8.9), Inches(0.3))
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Live session — skeleton overlay, form score, exercise demo panel"
    cp.font.size = Pt(11)
    cp.font.color.rgb = MUTED
    cp.alignment = PP_ALIGN.CENTER

    # Slide 8 — AI pipeline + pose visuals
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Pose Detection & Scoring")
    add_two_column(
        slide,
        "Detection",
        [
            "• MediaPipe PoseLandmarker (lite model)",
            "• 33 body landmarks per frame",
            "• Smoothed skeleton overlay on video",
            "• Runs in browser via WebAssembly",
        ],
        "Scoring",
        [
            "• Joint angles vs ideal template per exercise",
            "• Form score 0–100%",
            "• Classification: correct / incorrect / adjust",
            "• Coaching text + optional voice (Web Speech API)",
            "• Hold timer locks at 75%+ accuracy (yoga)",
        ],
    )
    add_image(slide, img["poses"], 1.2, 4.85, 7.6, 2.35)

    # Slide 9 — Exercises + pose reference images
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Exercise Catalog")
    add_image(slide, img["poses"], 0.55, 1.65, 8.9, 2.5)

    add_two_column(
        slide,
        "Fitness & strength",
        [
            "• Squat, Push-up, Lunge (rep counter)",
            "• Bicep curl, Shoulder press",
        ],
        "Yoga & office",
        [
            "• Mountain, Tree, Warrior, Cobra (hold)",
            "• Desk stretch, seated twist",
        ],
    )

    # Slide 10 — Database + dashboard
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Database & Dashboard")
    add_bullets(
        slide,
        [
            "profiles — user profile, onboarding, goal, BMI",
            "exercises — catalog with ideal_angles (JSONB)",
            "sessions, exercise_logs, progress_logs",
            "workout_plans + workout_plan_items",
            "user_streaks, user_equipment",
        ],
        width=4.5,
        height=4.5,
        size=15,
    )
    add_image(slide, img["dashboard"], 5.2, 1.75, 4.4, 2.48)

    # Slide 11 — APIs
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "REST API Endpoints")
    apis = [
        ("POST /api/auth/signup", "Create account (service role)"),
        ("GET/POST /api/sessions", "List / save workout sessions"),
        ("GET /api/exercises", "Exercise catalog"),
        ("GET /api/plan", "User workout plan"),
        ("GET /api/analytics/daily", "Daily activity report"),
        ("GET /api/analytics/weekly", "Weekly summary"),
    ]
    y = 1.9
    for route, desc in apis:
        rbox = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(3.2), Inches(0.35))
        rp = rbox.text_frame.paragraphs[0]
        rp.text = route
        rp.font.bold = True
        rp.font.size = Pt(13)
        rp.font.color.rgb = PRIMARY

        dbox = slide.shapes.add_textbox(Inches(4.0), Inches(y), Inches(5.5), Inches(0.35))
        dp = dbox.text_frame.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(13)
        dp.font.color.rgb = DARK
        y += 0.52

    # Slide 12 — Application pages + UI screenshots
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Application UI")
    add_image(slide, img["landing"], 0.55, 1.65, 4.3, 2.42)
    add_image(slide, img["session"], 5.15, 1.65, 4.3, 2.42)
    add_image(slide, img["dashboard"], 0.55, 4.25, 8.9, 2.42)
    labels = [
        (0.55, 4.05, "Landing"),
        (5.15, 4.05, "Live session"),
        (0.55, 6.65, "Dashboard"),
    ]
    for left, top, text in labels:
        lb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(4.3), Inches(0.25))
        lp = lb.text_frame.paragraphs[0]
        lp.text = text
        lp.font.size = Pt(11)
        lp.font.color.rgb = MUTED
        lp.alignment = PP_ALIGN.CENTER

    # Slide 13 — Setup & deploy
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Setup & Deployment")
    add_bullets(
        slide,
        [
            "Clone repo → npm run install:app → copy .env.local",
            "npm run setup:supabase — runs SQL migrations",
            "npm run dev → http://localhost:3000",
            "GitHub: github.com/BhanuTejapothuru01/yoga-pose-Tracker",
            "Production: yoga-tracker-kappa.vercel.app",
            "Supabase redirect URL must include the Vercel domain for auth",
        ],
        size=16,
    )

    # Slide 14 — Future work
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG)
    add_title_bar(slide, "Future Work")
    add_bullets(
        slide,
        [
            "Equipment detection from photos (YOLO model)",
            "More exercises and guided multi-pose flows",
            "Native mobile app with on-device models",
            "Social features and shared workout challenges",
            "Wearable heart-rate integration",
        ],
        size=17,
    )

    # Slide 15 — Thank you + logo
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PRIMARY)
    add_image(slide, img["logo"], 4.45, 1.2, 1.1)
    thanks = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(1.0))
    thp = thanks.text_frame.paragraphs[0]
    thp.text = "Thank You"
    thp.font.size = Pt(48)
    thp.font.bold = True
    thp.font.color.rgb = WHITE
    thp.alignment = PP_ALIGN.CENTER

    qa = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(8.4), Inches(0.6))
    qap = qa.text_frame.paragraphs[0]
    qap.text = "Questions?"
    qap.font.size = Pt(22)
    qap.font.color.rgb = ACCENT
    qap.alignment = PP_ALIGN.CENTER

    footer = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(8.4), Inches(0.8))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "YogaTracker — yoga-tracker-kappa.vercel.app"
    fp.font.size = Pt(14)
    fp.font.color.rgb = WHITE
    fp.alignment = PP_ALIGN.CENTER

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Presentation saved to: {path}")
